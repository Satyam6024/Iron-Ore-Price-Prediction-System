from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = PROJECT_ROOT / "presentation"
FIGURE_DIR = PRESENTATION_DIR / "figures"
OUTPUT_PATH = PRESENTATION_DIR / "Iron_Ore_Project_Presentation.pptx"


SLIDES = [
    {
        "title": "End-to-End Iron Ore Price Forecasting and MLOps Workflow",
        "bullets": [
            "Problem: forecast iron ore prices under regime changes and market turbulence",
            "Scope: production ML pipeline plus research-aligned benchmarking",
            "Stack: DVC, MLflow, FastAPI, Docker, Jupyter notebooks",
        ],
        "image": "09_mlop_flow.png",
    },
    {
        "title": "Iron Ore Market Dynamics",
        "bullets": [
            "2,264 trading-day records from 2016-01-04 to 2024-12-31",
            "Core variables: Price, Open, High, Low, Volume, Change %",
            "Series shows trends, reversals, and volatility bursts",
            "Static forecasting assumptions are not reliable here",
        ],
        "image": "01_market_dynamics.png",
    },
    {
        "title": "Data Ingestion, Cleaning, and EDA",
        "bullets": [
            "Raw CSV is ingested into a DVC-tracked artifact",
            "Cleaning standardizes schema, parses dates, and converts numeric fields",
            "Volume and percentage columns are normalized for modeling",
            "EDA checks seasonality, volatility, and feature relationships",
        ],
        "image": "03_seasonality.png",
    },
    {
        "title": "Volatility and Conditional Variance",
        "bullets": [
            "Returns are unstable across time and cluster during turbulent periods",
            "Research workflow tracks a 14-day rolling volatility proxy",
            "A 75th percentile threshold separates high- and low-volatility regimes",
            "EWMA variance error quantifies conditional-variance mismatch",
        ],
        "image": "02_volatility_regime.png",
    },
    {
        "title": "Linear Models Under Turbulence",
        "bullets": [
            "After correcting to rolling one-step evaluation, ARIMA and Holt-Winters remain strong baselines",
            "ARIMA RMSE: 1.862 and Holt-Winters RMSE: 1.860 on the research benchmark",
            "Their high-volatility RMSE still stays above 3.18, where turbulence exposes the limits of linear structure",
            "This supports residual hybridization rather than discarding statistical models entirely",
        ],
        "image": "06_research_benchmark.png",
    },
    {
        "title": "Machine Learning Approaches",
        "bullets": [
            "SVR formulates nonlinear forecasting as a convex optimization problem in a kernel-induced high-dimensional feature space",
            "SVR accuracy depends on epsilon-insensitive regularization, kernel choice, hyperparameter calibration, and lag-feature design",
            "Random Forest and Gradient Boosting reduce variance via feature subsampling, bootstrap aggregation, and sequential tree refinement",
            "Core limitation: these models do not natively preserve ordered temporal memory without explicit lag engineering",
            "SVR is the strongest pure ML model at RMSE 1.9082, but Hybrid ARIMA + LSTM improves further to 1.8576",
        ],
        "image": "13_ml_model_comparison.png",
    },
    {
        "title": "Deep Learning for Sequence Modeling",
        "bullets": [
            "LSTM models internalize sequential state dynamics instead of relying only on handcrafted lags",
            "In this project, the residual LSTM stopped after 26 epochs and reached its best validation loss at epoch 14",
            "The recurrent block models ARIMA residuals rather than raw price directly",
            "This turns the paper's ARIMA-LSTM idea into a real benchmarked experiment in the repo",
        ],
        "image": "14_arima_lstm_history.png",
    },
    {
        "title": "The Hybrid Ensemble Philosophy",
        "bullets": [
            "Hybrid models assume linear and nonlinear components should be learned separately",
            "ARIMA provides the rolling one-step linear forecast and LSTM models the nonlinear residual structure",
            "Hybrid ARIMA + LSTM is the best research model with RMSE 1.8576 and R2 0.9760",
            "High-volatility RMSE falls to 3.1817, slightly better than ARIMA, Holt-Winters, and SVR",
            "The gain is small but consistent, which is typical of residual-correction hybrids",
        ],
        "image": "07_research_prediction_traces.png",
    },
    {
        "title": "Sequential Decomposition (ARIMA-LSTM) as the Next Step",
        "bullets": [
            "Step 1: fit ARIMA on the raw series for the linear component",
            "Step 2: compute residuals from ARIMA forecasts",
            "Step 3: train LSTM on residual sequences for nonlinear structure",
            "Step 4: combine both forecasts into the final signal",
        ],
        "image": "08_sequential_decomposition.png",
    },
    {
        "title": "Performance Evaluation",
        "bullets": [
            "Evaluation is not limited to RMSE",
            "Production reports RMSE, MAE, MAPE, and R2",
            "Research adds regime RMSE, rolling RMSE, procurement-cost proxy, and EWMA variance error",
            "Best research result is Hybrid ARIMA + LSTM with RMSE 1.8576 and procurement-cost proxy 2.22M",
            "This makes evaluation closer to business and risk reality",
        ],
        "image": "10_production_fit.png",
    },
    {
        "title": "Identified Research Gaps",
        "bullets": [
            "Conditional variance is diagnosed but not yet modeled explicitly with GARCH-family methods",
            "Structural-break detection is not yet part of the pipeline",
            "ARIMA-LSTM is benchmarked, but not yet wired into production training, registry, or deployment",
            "Exogenous macro or event variables are still missing",
        ],
        "image": "04_correlation_heatmap.png",
    },
    {
        "title": "MLOps Workflow and Deployment",
        "bullets": [
            "DVC orchestrates ingestion, cleaning, feature engineering, selection, training, and deployment stages",
            "MLflow tracks experiments and manages model lifecycle",
            "Promotion gates move models across testing, staging, and production",
            "FastAPI serves predictions and Docker packages the stack",
        ],
        "image": "09_mlop_flow.png",
    },
    {
        "title": "Conclusion",
        "bullets": [
            "The project combines classical forecasting, machine learning, and MLOps in one workflow",
            "Production pipeline delivers a strong deployable Ridge model on selected engineered features",
            "Research pipeline now includes a real ARIMA + LSTM residual hybrid with volatility-aware diagnostics",
            "The strongest next step is richer variance modeling and structural-break detection around that hybrid core",
        ],
        "image": "01_market_dynamics.png",
    },
]


def add_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7))
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(24)
    p.font.bold = True


def add_bullets(slide, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4.6), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    font_size = 17 if len(bullets) <= 4 else 15.5
    space_after = 10 if len(bullets) <= 4 else 7
    first = True
    for bullet in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.space_after = Pt(space_after)
        first = False


def add_picture(slide, image_name: str) -> None:
    image_path = FIGURE_DIR / image_name
    slide.shapes.add_picture(str(image_path), Inches(5.45), Inches(1.2), width=Inches(7.4))


def add_images(slide, images: list[dict]) -> None:
    for image in images:
        image_path = FIGURE_DIR / image["name"]
        slide.shapes.add_picture(
            str(image_path),
            Inches(image["left"]),
            Inches(image["top"]),
            width=Inches(image["width"]),
        )


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    for slide_cfg in SLIDES:
        slide = prs.slides.add_slide(blank)
        add_title(slide, slide_cfg["title"])
        add_bullets(slide, slide_cfg["bullets"])
        if "images" in slide_cfg:
            add_images(slide, slide_cfg["images"])
        else:
            add_picture(slide, slide_cfg["image"])

    prs.save(OUTPUT_PATH)
    print(f"[presentation] built {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
